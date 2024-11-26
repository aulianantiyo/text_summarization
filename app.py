import streamlit as st
import nltk
from summa import summarizer
from transformers import pipeline
from sklearn.feature_extraction.text import TfidfVectorizer
from nltk.tokenize import sent_tokenize
import PyPDF2
import io
import docx
from pptx import Presentation

# Download required NLTK data
nltk.download('punkt')

def extract_text_from_pdf(file_bytes):
    """Extract text from PDF file"""
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file_bytes):
    """Extract text from DOCX file"""
    doc = docx.Document(io.BytesIO(file_bytes))
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

def extract_text_from_pptx(file_bytes):
    """Extract text from PPTX file"""
    presentation = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text.append(paragraph.text)
    return '\n'.join(text)

def extractive_summary(text: str, ratio: float = 0.3) -> str:
    """Generate extractive summary using TextRank algorithm"""
    return summarizer.summarize(text, ratio=ratio)

def tfidf_summary(text: str, num_sentences: int = 5) -> str:
    """Generate extractive summary using TF-IDF"""
    sentences = sent_tokenize(text)
    if len(sentences) <= num_sentences:
        return text
        
    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform(sentences)
    scores = tfidf_matrix.sum(axis=1).tolist()
    ranked_sentences = [x for _, x in sorted(zip(scores, sentences), reverse=True)]
    return ' '.join(ranked_sentences[:num_sentences])

@st.cache_resource
def load_summarizer():
    """Load and cache the BART summarizer"""
    return pipeline("summarization", model="facebook/bart-large-cnn")

def abstractive_summary(text: str) -> str:
    """Generate abstractive summary using BART model"""
    summarizer = load_summarizer()
    # Split text into chunks if it's too long
    max_chunk_length = 1024
    chunks = [text[i:i + max_chunk_length] for i in range(0, len(text), max_chunk_length)]
    
    summaries = []
    for chunk in chunks:
        summary = summarizer(chunk, max_length=130, min_length=30, do_sample=False)
        summaries.append(summary[0]['summary_text'])
    
    return ' '.join(summaries)

def main():
    st.set_page_config(
        page_title="Summer",
        page_icon="☀️📄",
        layout="wide"
    )

    # Header
    st.title("☀️📄 Summary Retrieval Assistant")
    st.markdown("Upload your document and get an AI-powered summary")

    # Sidebar for settings
    with st.sidebar:
        st.header("Options")
        method = st.selectbox(
            "Summarization Method",
            ['extractive', 'abstractive', 'tf-idf'],
            help="Choose the summarization technique"
        )
    
        if method == 'extractive':
            ratio = st.slider(
                "Summary Length Ratio",
                min_value=0.1,
                max_value=0.5,
                value=0.3,
                step=0.1,
                help="Ratio of original text to keep in summary"
            )

    # Main content area
    uploaded_file = st.file_uploader(
        "Choose a file",
        type=['txt', 'pdf', 'docx', 'pptx'],
        help="Upload a text, PDF, Word, or PowerPoint document"
    )

    if uploaded_file is not None:
        # Show file details
        col1, col2 = st.columns(2)
        with col1:
            st.write(f"Filename: {uploaded_file.name}")
        with col2:
            st.write(f"File size: {uploaded_file.size / 1024:.2f} KB")

        # Process button
        if st.button("Generate Summary", type="primary"):
            with st.spinner('Processing your document...'):
                try:
                    # Read file content based on file type
                    if uploaded_file.type == "application/pdf":
                        text = extract_text_from_pdf(uploaded_file.read())
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                        text = extract_text_from_docx(uploaded_file.read())
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                        text = extract_text_from_pptx(uploaded_file.read())
                    else:
                        text = uploaded_file.read().decode()

                    # Generate summary based on selected method
                    if method == "extractive":
                        summary = extractive_summary(text, ratio)
                    elif method == "tf-idf":
                        summary = tfidf_summary(text)
                    else:
                        summary = abstractive_summary(text)

                    # Display results
                    st.header("Results")
                    
                    # Display statistics
                    col1, col2, col3 = st.columns(3)
                    col1.metric("Original Text Length", f"{len(text)} chars")
                    col2.metric("Summary Length", f"{len(summary)} chars")
                    col3.metric("Compression Ratio", f"{(len(summary)/len(text))*100:.1f}%")

                    # Display original text and summary in tabs
                    tab1, tab2 = st.tabs(["Summary ", "Original Text"])
                    
                    with tab1:
                        st.subheader("Summary")
                        st.write(summary)
                        
                        # Add download button for summary
                        st.download_button(
                            label="Download Summary",
                            data=summary,
                            file_name=f"summary_{uploaded_file.name}.txt",
                            mime="text/plain"
                        )
                    
                    with tab2:
                        st.subheader("Original Text")
                        st.write(text)

                except Exception as e:
                    st.error(f"An error occurred: {str(e)}")
                    st.error("Please make sure the uploaded file is valid and try again.")

    # Add information about the methods
    with st.expander("About the Summarization Methods"):
        st.markdown("""
        ### Available Methods:
        1. **Extractive Summarization**
           - Uses TextRank algorithm to identify key sentences
           - Maintains original wording
           - Allows control over summary length
        
        2. **Abstractive Summarization**
           - Uses BART neural network to generate new text
           - Can rephrase and combine information
           - May produce more natural summaries
        
        3. **TF-IDF Summarization**
           - Uses term frequency to identify important sentences
           - Good for technical documents
           - Maintains original wording
        """)

if __name__ == "__main__":
    main()
